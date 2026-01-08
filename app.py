"""
企业碳排放计算器 - 专业版
基于GHG Protocol和IPCC 2006标准
"""

import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment

# 页面配置
st.set_page_config(
    page_title="企业碳排放计算器 - 专业版",
    page_icon="🌱",
    layout="wide"
)

# CSS样式
st.markdown("""
<style>
    .main-header {font-size: 2.5rem; font-weight: 700; color: #1e3a8a; text-align: center; padding: 1rem 0;}
    .sub-header {font-size: 1.2rem; color: #64748b; text-align: center; padding-bottom: 2rem;}
</style>
""", unsafe_allow_html=True)

# 初始化排放因子库
if 'emission_factors' not in st.session_state:
    st.session_state.emission_factors = {
        "固定燃烧-天然气": {"factor": 2.1622, "unit": "kgCO2/m3", "ghg_type": "CO2"},
        "固定燃烧-煤炭": {"factor": 2.38, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "固定燃烧-柴油": {"factor": 3.0959, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "固定燃烧-汽油": {"factor": 2.9251, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "移动燃烧-汽油": {"factor": 2.9251, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "移动燃烧-柴油": {"factor": 3.0959, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "工艺排放-丙烷": {"factor": 2.9761, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "工艺排放-二氧化碳": {"factor": 1.0, "unit": "kgCO2/kg", "ghg_type": "CO2"},
        "无组织排放-R410A": {"factor": 2088, "unit": "kgCO2e/kg", "ghg_type": "HFCs"},
        "无组织排放-R32": {"factor": 675, "unit": "kgCO2e/kg", "ghg_type": "HFCs"},
        "无组织排放-甲烷(化粪池)": {"factor": 22.4, "unit": "kgCO2e/kgBOD", "ghg_type": "CH4"},
        "外购电力-全国平均": {"factor": 0.5703, "unit": "kgCO2/kWh", "ghg_type": "CO2"},
        "外购电力-华北区域": {"factor": 0.8843, "unit": "kgCO2/kWh", "ghg_type": "CO2"},
        "外购电力-华东区域": {"factor": 0.7035, "unit": "kgCO2/kWh", "ghg_type": "CO2"},
        "外购热力-蒸汽": {"factor": 110, "unit": "kgCO2/GJ", "ghg_type": "CO2"},
    }

if 'matched_data' not in st.session_state:
    st.session_state.matched_data = None
if 'calculation_done' not in st.session_state:
    st.session_state.calculation_done = False

# 侧边栏
with st.sidebar:
    st.title("🔧 排放因子管理")
    st.markdown("---")
    
    total_factors = len(st.session_state.emission_factors)
    st.metric("因子总数", total_factors)
    
    with st.expander("📚 查看因子库"):
        factor_df = pd.DataFrame([
            {'排放源': k, '排放因子': v['factor'], '单位': v['unit'], '气体': v['ghg_type']}
            for k, v in st.session_state.emission_factors.items()
        ])
        st.dataframe(factor_df, use_container_width=True)
    
    st.subheader("➕ 添加排放因子")
    with st.form("add_factor"):
        new_name = st.text_input("排放源名称", placeholder="例：固定燃烧-生物质")
        new_factor = st.number_input("排放因子", min_value=0.0, step=0.01, format="%.4f")
        new_unit = st.text_input("单位", placeholder="kgCO2/kg")
        new_ghg = st.selectbox("温室气体", ["CO2", "CH4", "N2O", "HFCs"])
        
        if st.form_submit_button("✅ 添加"):
            if new_name and new_factor > 0 and new_unit:
                st.session_state.emission_factors[new_name] = {
                    "factor": new_factor, "unit": new_unit, "ghg_type": new_ghg
                }
                st.success(f"✅ 已添加: {new_name}")
                st.rerun()

# 主界面
st.markdown('<p class="main-header">🌱 企业碳排放计算器 - 专业版</p>', unsafe_allow_html=True)
st.markdown('<p class="sub-header">基于 GHG Protocol 和 IPCC 2006 标准</p>', unsafe_allow_html=True)

# 创建模板
def create_template():
    data = {
        '类别': ['范围一：直接温室气体排放']*4 + ['范围二：间接温室气体排放']*2,
        '子类别': ['1.1 固定燃烧', '1.2 移动燃烧', '1.3 工艺排放', '1.4 无组织排放', '2.1 外购电力', '2.2 外购热力'],
        '排放源': ['天然气', '汽油', '丙烷', 'R410A', '外购市政电', '蒸汽'],
        '设施/过程': ['燃气锅炉', '公务车', '焊接', '空调', '用电', '供暖设备'],
        '活动数据': [1239138, 11010, 792, 3.15, 1500000, 500],
        '计量单位': ['m³', 'kg', 'kg', 'kg', 'kWh', 'GJ']
    }
    df = pd.DataFrame(data)
    
    output = BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df.to_excel(writer, index=False, sheet_name='活动数据')
        ws = writer.sheets['活动数据']
        for i, col in enumerate(['A', 'B', 'C', 'D', 'E', 'F'], 1):
            ws.column_dimensions[col].width = 25
        for cell in ws[1]:
            cell.fill = PatternFill(start_color='4472C4', end_color='4472C4', fill_type='solid')
            cell.font = Font(color='FFFFFF', bold=True)
    return output.getvalue()

# 步骤1: 下载模板
st.subheader("📥 步骤1: 下载活动数据模板")
col1, col2 = st.columns([3, 1])
with col1:
    st.info("📌 包含范围一（固定燃烧、移动燃烧、工艺排放、无组织排放）+ 范围二（外购电力、热力）")
with col2:
    st.download_button("📄 下载模板", create_template(), 
                      "碳排放数据模板.xlsx",
                      "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                      use_container_width=True)

st.markdown("---")

# 步骤2: 上传数据
st.subheader("📤 步骤2: 上传活动数据")
uploaded_file = st.file_uploader("上传Excel文件", type=['xlsx', 'xls'])

if uploaded_file:
    try:
        df = pd.read_excel(uploaded_file)
        required_cols = ['类别', '子类别', '排放源', '设施/过程', '活动数据', '计量单位']
        
        if not all(col in df.columns for col in required_cols):
            st.error(f"❌ 文件格式不正确！必需列：{', '.join(required_cols)}")
        else:
            st.success("✅ 文件上传成功！")
            
            with st.expander("📊 查看上传数据", expanded=True):
                st.dataframe(df, use_container_width=True)
            
            st.markdown("---")
            st.subheader("🔍 步骤3: 排放因子智能匹配")
            
            if st.button("🚀 开始匹配排放因子", type="primary", use_container_width=True):
                result_df = df.copy()
                result_df['建议排放源类型'] = None
                result_df['排放因子'] = None
                result_df['因子单位'] = None
                result_df['温室气体类型'] = None
                result_df['匹配状态'] = None
                
                for idx, row in result_df.iterrows():
                    subcat = str(row['子类别'])
                    source = str(row['排放源'])
                    
                    if '1.1' in subcat:
                        key = f"固定燃烧-{source}"
                    elif '1.2' in subcat:
                        key = f"移动燃烧-{source}"
                    elif '1.3' in subcat:
                        key = f"工艺排放-{source}"
                    elif '1.4' in subcat:
                        key = f"无组织排放-{source}"
                    elif '2.1' in subcat:
                        key = "外购电力-全国平均" if "电" in source else f"外购电力-{source}"
                    elif '2.2' in subcat:
                        key = f"外购热力-{source}"
                    else:
                        key = None
                    
                    if key and key in st.session_state.emission_factors:
                        info = st.session_state.emission_factors[key]
                        result_df.at[idx, '建议排放源类型'] = key
                        result_df.at[idx, '排放因子'] = info['factor']
                        result_df.at[idx, '因子单位'] = info['unit']
                        result_df.at[idx, '温室气体类型'] = info['ghg_type']
                        result_df.at[idx, '匹配状态'] = '✅ 已匹配'
                    else:
                        result_df.at[idx, '建议排放源类型'] = key or "未识别"
                        result_df.at[idx, '排放因子'] = 0
                        result_df.at[idx, '因子单位'] = '待补充'
                        result_df.at[idx, '温室气体类型'] = 'CO2'
                        result_df.at[idx, '匹配状态'] = '❌ 未匹配'
                
                st.session_state.matched_data = result_df
                st.success("✅ 匹配完成！")
            
            if st.session_state.matched_data is not None:
                st.markdown("#### 📋 匹配结果")
                matched_df = st.session_state.matched_data
                
                col1, col2, col3 = st.columns(3)
                total = len(matched_df)
                matched = len(matched_df[matched_df['匹配状态'] == '✅ 已匹配'])
                col1.metric("总数", total)
                col2.metric("已匹配", matched)
                col3.metric("未匹配", total - matched)
                
                edited_df = st.data_editor(matched_df, use_container_width=True, height=400,
                    column_config={
                        "建议排放源类型": st.column_config.SelectboxColumn(
                            "建议排放源类型",
                            options=list(st.session_state.emission_factors.keys())
                        ),
                        "排放因子": st.column_config.NumberColumn("排放因子", format="%.4f")
                    },
                    disabled=['类别', '子类别', '排放源', '设施/过程', '活动数据', '计量单位']
                )
                
                st.session_state.matched_data = edited_df
                
                if st.button("✅ 确认匹配，开始计算", type="primary", use_container_width=True):
                    st.session_state.calculation_done = True
                    st.rerun()
    
    except Exception as e:
        st.error(f"❌ 文件读取失败: {str(e)}")

# 步骤4: 计算和可视化
if st.session_state.calculation_done and st.session_state.matched_data is not None:
    st.markdown("---")
    st.subheader("📊 步骤4: 排放计算结果")
    
    calc_df = st.session_state.matched_data.copy()
    calc_df['排放量(kgCO2e)'] = calc_df['活动数据'] * calc_df['排放因子']
    calc_df['排放量(tCO2e)'] = calc_df['排放量(kgCO2e)'] / 1000
    calc_df['范围'] = calc_df['类别'].apply(lambda x: '范围一' if '直接' in x else '范围二')
    
    total_emission = calc_df['排放量(tCO2e)'].sum()
    scope_summary = calc_df.groupby('范围')['排放量(tCO2e)'].sum()
    scope1 = scope_summary.get('范围一', 0)
    scope2 = scope_summary.get('范围二', 0)
    
    # 汇总卡片
    st.markdown("### 📈 排放汇总")
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown(f"""
        <div style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); 
                    padding: 1.5rem; border-radius: 10px; color: white;'>
            <h3 style='margin:0;'>范围一：直接排放</h3>
            <h2 style='margin:0.5rem 0 0 0; font-size: 2.5rem;'>{scope1:.2f}</h2>
            <p style='margin:0;'>tCO₂e</p>
        </div>
        """, unsafe_allow_html=True)
    
    with col2:
        st.markdown(f"""
        <div style='background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); 
                    padding: 1.5rem; border-radius: 10px; color: white;'>
            <h3 style='margin:0;'>范围二：间接排放</h3>
            <h2 style='margin:0.5rem 0 0 0; font-size: 2.5rem;'>{scope2:.2f}</h2>
            <p style='margin:0;'>tCO₂e</p>
        </div>
        """, unsafe_allow_html=True)
    
    with col3:
        st.markdown(f"""
        <div style='background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%); 
                    padding: 1.5rem; border-radius: 10px; color: white;'>
            <h3 style='margin:0;'>排放总量</h3>
            <h2 style='margin:0.5rem 0 0 0; font-size: 2.5rem;'>{total_emission:.2f}</h2>
            <p style='margin:0;'>tCO₂e</p>
        </div>
        """, unsafe_allow_html=True)
    
    st.markdown("---")
    
    # 可视化
    tab1, tab2, tab3 = st.tabs(["📊 温室气体分析", "🔥 排放源分析", "📑 数据导出"])
    
    with tab1:
        col1, col2 = st.columns(2)
        
        with col1:
            ghg_summary = calc_df.groupby('温室气体类型')['排放量(tCO2e)'].sum().reset_index()
            fig = px.pie(ghg_summary, values='排放量(tCO2e)', names='温室气体类型',
                        title='温室气体排放占比', hole=0.4)
            st.plotly_chart(fig, use_container_width=True)
        
        with col2:
            fig2 = px.bar(ghg_summary, x='温室气体类型', y='排放量(tCO2e)',
                         title='各温室气体排放量')
            st.plotly_chart(fig2, use_container_width=True)
    
    with tab2:
        col1, col2 = st.columns(2)
        
        with col1:
            scope_df = pd.DataFrame({'范围': ['范围一', '范围二'], '排放量': [scope1, scope2]})
            fig3 = px.pie(scope_df, values='排放量', names='范围',
                         title='范围一 vs 范围二', hole=0.4)
            st.plotly_chart(fig3, use_container_width=True)
        
        with col2:
            subcat = calc_df.groupby('子类别')['排放量(tCO2e)'].sum().reset_index()
            fig4 = px.bar(subcat, x='子类别', y='排放量(tCO2e)',
                         title='各子类别排放量')
            st.plotly_chart(fig4, use_container_width=True)
    
    with tab3:
        col1, col2 = st.columns(2)
        
        with col1:
            def export_excel():
                output = BytesIO()
                with pd.ExcelWriter(output, engine='openpyxl') as writer:
                    calc_df.to_excel(writer, sheet_name='详细数据', index=False)
                    summary = pd.DataFrame({
                        '指标': ['范围一', '范围二', '总量'],
                        '排放量(tCO2e)': [scope1, scope2, total_emission]
                    })
                    summary.to_excel(writer, sheet_name='汇总', index=False)
                return output.getvalue()
            
            st.download_button("📥 下载Excel报告", export_excel(),
                              f"碳排放报告_{pd.Timestamp.now().strftime('%Y%m%d')}.xlsx",
                              use_container_width=True)
        
        with col2:
            def create_ppt():
                prs = Presentation()
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
                
                # 封面
                slide1 = prs.slides.add_slide(prs.slide_layouts[6])
                slide1.background.fill.solid()
                slide1.background.fill.fore_color.rgb = RGBColor(30, 58, 138)
                
                title = slide1.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(1.5))
                tf = title.text_frame
                tf.text = "企业碳排放计算报告"
                tf.paragraphs[0].font.size = Pt(54)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                subtitle = slide1.shapes.add_textbox(Inches(2), Inches(5), Inches(12), Inches(1))
                stf = subtitle.text_frame
                stf.text = f"总排放量: {total_emission:.2f} tCO₂e"
                stf.paragraphs[0].font.size = Pt(36)
                stf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                stf.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 数据页
                slide2 = prs.slides.add_slide(prs.slide_layouts[6])
                title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(0.8))
                tf2 = title2.text_frame
                tf2.text = "排放汇总分析"
                tf2.paragraphs[0].font.size = Pt(40)
                tf2.paragraphs[0].font.bold = True
                
                table = slide2.shapes.add_table(3, 3, Inches(2), Inches(2), Inches(12), Inches(2.5)).table
                headers = ['范围', '排放量(tCO₂e)', '占比']
                for i, h in enumerate(headers):
                    table.cell(0, i).text = h
                
                table.cell(1, 0).text = "范围一"
                table.cell(1, 1).text = f"{scope1:.2f}"
                table.cell(1, 2).text = f"{scope1/total_emission*100:.1f}%"
                
                table.cell(2, 0).text = "范围二"
                table.cell(2, 1).text = f"{scope2:.2f}"
                table.cell(2, 2).text = f"{scope2/total_emission*100:.1f}%"
                
                output = BytesIO()
                prs.save(output)
                return output.getvalue()
            
            st.download_button("📥 生成PPT报告", create_ppt(),
                              f"碳排放报告_{pd.Timestamp.now().strftime('%Y%m%d')}.pptx",
                              use_container_width=True)

st.markdown("---")
st.markdown("""
<div style='text-align: center; color: #666; padding: 2rem;'>
    <p>🌱 企业碳排放计算器 - 专业版 v2.0</p>
    <p>基于 GHG Protocol 和 IPCC 2006 标准</p>
</div>
""", unsafe_allow_html=True)